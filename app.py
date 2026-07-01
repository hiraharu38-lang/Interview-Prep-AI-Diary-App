import streamlit as st
from google import genai
from pptx import Presentation
from pypdf import PdfReader
from docx import Document
import json
import time
import io
import os
import copy

st.set_page_config(page_title="Slide Generator AI PRO", page_icon="🎨", layout="wide")

# ============================================================
# 初期設定
# ============================================================
SELECT_MODEL = "gemini-2.5-flash"
api_key = st.secrets["GEMINI_API_KEY"]
client = genai.Client(api_key=api_key)
IMAGE_DIR = "./images/"
TEMPLATE_PATH = "./master_template.pptx"

# テンプレート内の「お手本スライド」のインデックス（0始まり）と役割の対応
MODEL_INDEX = {
    "title": 0,            # Title Slide with Image
    "agenda": 1,           # Agenda
    "intro": 2,            # Introduction（タイトル＋本文＋画像）
    "chapter": 3,          # Section Header（章見出し／強調テキスト）
    "two_content": 11,     # Two Content（2カラム）
    "three_content": 12,   # Three Content（3カラム＋アイコン画像）
    "quote": 5,            # Quote
    "summary": 13,         # Summary
    "thank_you": 14,       # Thank you
}

# ============================================================
# リトライ付きGemini呼び出し
# ============================================================
def generate_with_retry(prompt, max_retries=6):
    for i in range(max_retries):
        try:
            response = client.models.generate_content(
                model=SELECT_MODEL,
                contents=prompt,
                config={"response_mime_type": "application/json"},
            )
            return response.text
        except Exception as e:
            if i < max_retries - 1:
                time.sleep(2.0)
                continue
            raise e

# ============================================================
# ファイルからテキスト抽出（PDF / Word）
# ============================================================
def extract_text(uploaded_file):
    if uploaded_file.type == "application/pdf":
        reader = PdfReader(uploaded_file)
        return "".join([(page.extract_text() or "") for page in reader.pages])
    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = Document(uploaded_file)
        return "\n".join([para.text for para in doc.paragraphs])
    return ""

# ============================================================
# テンプレートのスライドを「同じプレゼンテーション内で」複製する
# （同じスライドマスター／レイアウトを参照するので、背景や装飾図形がそのまま維持される）
# ============================================================
def duplicate_slide(prs, index):
    source = prs.slides[index]
    layout = source.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # add_slide() がレイアウトから自動生成した空のプレースホルダーを一旦すべて削除
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)

    # 元スライドの図形（プレースホルダー・装飾図形すべて）をXMLごと複製
    for shp in source.shapes:
        new_el = copy.deepcopy(shp._element)
        new_slide.shapes._spTree.append(new_el)

    return new_slide

def get_ph(slide, idx):
    try:
        return slide.placeholders[idx]
    except KeyError:
        return None

def set_text_preserve_format(ph, text):
    """プレースホルダーの先頭ランの書式（フォント・色・サイズ）を保ったままテキストだけ差し替える"""
    if ph is None or text is None:
        return
    tf = ph.text_frame
    paragraphs = list(tf.paragraphs)
    if not paragraphs:
        tf.text = text
        return
    p0 = paragraphs[0]
    if p0.runs:
        p0.runs[0].text = text
        for r in list(p0.runs[1:]):
            r._r.getparent().remove(r._r)
    else:
        p0.text = text
    # 2段落目以降（テンプレートのダミー複数行）は削除
    for p in paragraphs[1:]:
        p._p.getparent().remove(p._p)

def set_multiline_preserve_format(ph, lines):
    """1つ目の段落の書式をそのままコピーして複数行（箇条書き）を差し込む"""
    if ph is None:
        return
    lines = [l for l in (lines or []) if l]
    if not lines:
        return
    tf = ph.text_frame
    paragraphs = list(tf.paragraphs)
    if not paragraphs:
        tf.text = lines[0]
        return
    base_el = paragraphs[0]._p
    parent = base_el.getparent()

    # まず1行目を既存の段落に反映
    set_text_preserve_format(ph, lines[0])

    # 2行目以降は1行目の段落XMLを複製して追加
    base_copy_source = copy.deepcopy(base_el)
    for line in lines[1:]:
        new_p = copy.deepcopy(base_copy_source)
        # 複製した段落内のテキストランをすべて1つにまとめて差し替え
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        r_elements = new_p.findall("a:r", ns)
        if r_elements:
            t_el = r_elements[0].find("a:t", ns)
            if t_el is not None:
                t_el.text = line
            for extra_r in r_elements[1:]:
                new_p.remove(extra_r)
        else:
            # ランが無い場合は素直にテキストを追加
            from pptx.oxml.ns import qn
            r = new_p.makeelement(qn("a:r"), {})
            t = r.makeelement(qn("a:t"), {})
            t.text = line
            r.append(t)
            new_p.append(r)
        parent.append(new_p)

def set_picture(slide, idx, image_name):
    """PICTUREプレースホルダーに画像を差し込む（テンプレートの枠どおりに自動クロップされる）"""
    if not image_name or image_name == "none":
        return False
    path = os.path.join(IMAGE_DIR, f"{image_name}.png")
    if not os.path.exists(path):
        return False
    ph = get_ph(slide, idx)
    if ph is None:
        return False
    try:
        ph.insert_picture(path)
        return True
    except Exception:
        return False

def set_footer_everywhere(prs, presentation_title, start_index):
    """新規生成したスライド全てのフッターに、プレゼン全体のタイトルを反映"""
    for slide in list(prs.slides)[start_index:]:
        for shp in slide.placeholders:
            try:
                if shp.placeholder_format.type is not None and "FOOTER" in str(shp.placeholder_format.type):
                    set_text_preserve_format(shp, presentation_title)
            except Exception:
                continue

# ============================================================
# レイアウトごとの差し込み処理
# ============================================================
def fill_title(slide, data):
    set_text_preserve_format(get_ph(slide, 0), data.get("title", ""))
    set_text_preserve_format(get_ph(slide, 28), data.get("subtitle", ""))
    set_picture(slide, 47, data.get("image_name"))

def fill_chapter(slide, data):
    set_text_preserve_format(get_ph(slide, 0), data.get("title", ""))
    set_text_preserve_format(get_ph(slide, 28), data.get("eyebrow", ""))
    set_picture(slide, 47, data.get("image_name"))

def fill_intro(slide, data):
    set_text_preserve_format(get_ph(slide, 0), data.get("title", ""))
    set_text_preserve_format(get_ph(slide, 28), data.get("body", ""))
    set_picture(slide, 51, data.get("image_name"))

def fill_agenda(slide, data):
    set_text_preserve_format(get_ph(slide, 0), data.get("title", "Agenda"))
    items = data.get("items", [])[:5]
    idx_list = [28, 29, 30, 31, 32]
    for i, idx in enumerate(idx_list):
        text = items[i] if i < len(items) else ""
        set_text_preserve_format(get_ph(slide, idx), text)

def fill_two_content(slide, data):
    set_text_preserve_format(get_ph(slide, 0), data.get("title", ""))
    cols = data.get("columns", [])
    pairs = [(27, 28), (52, 53)]
    for i, (h_idx, b_idx) in enumerate(pairs):
        col = cols[i] if i < len(cols) else {}
        set_text_preserve_format(get_ph(slide, h_idx), col.get("heading", ""))
        set_text_preserve_format(get_ph(slide, b_idx), col.get("body", ""))
    set_picture(slide, 51, data.get("image_name"))

def fill_three_content(slide, data):
    set_text_preserve_format(get_ph(slide, 0), data.get("title", ""))
    items = data.get("items", [])
    groups = [(27, 28, 36), (29, 34, 37), (31, 35, 38)]
    for i, (h_idx, b_idx, pic_idx) in enumerate(groups):
        item = items[i] if i < len(items) else {}
        set_text_preserve_format(get_ph(slide, h_idx), item.get("heading", ""))
        set_text_preserve_format(get_ph(slide, b_idx), item.get("body", ""))
        set_picture(slide, pic_idx, item.get("image_name"))

def fill_quote(slide, data):
    set_text_preserve_format(get_ph(slide, 0), data.get("quote", ""))
    set_text_preserve_format(get_ph(slide, 29), data.get("author", ""))

def fill_summary(slide, data):
    set_text_preserve_format(get_ph(slide, 0), data.get("title", "まとめ"))
    set_multiline_preserve_format(get_ph(slide, 28), data.get("bullets", []))
    set_picture(slide, 48, data.get("image_name"))

def fill_thank_you(slide, data):
    set_multiline_preserve_format(get_ph(slide, 27), data.get("contact_lines", []))

FILLERS = {
    "title": fill_title,
    "chapter": fill_chapter,
    "intro": fill_intro,
    "agenda": fill_agenda,
    "two_content": fill_two_content,
    "three_content": fill_three_content,
    "quote": fill_quote,
    "summary": fill_summary,
    "thank_you": fill_thank_you,
}

# ============================================================
# PPTX生成本体（テンプレート方式）
# ============================================================
def create_pptx(payload):
    if not os.path.exists(TEMPLATE_PATH):
        raise FileNotFoundError("master_template.pptx が見つかりません。app.py と同じフォルダに置いてください。")

    prs = Presentation(TEMPLATE_PATH)
    original_sld_ids = list(prs.slides._sldIdLst)  # 元々入っている15枚（お手本）のID一覧
    presentation_title = payload.get("presentation_title", "")
    slides_data = payload.get("slides", [])

    start_index = len(prs.slides)  # ここから新規スライドが追加されていく

    for data in slides_data:
        layout_key = data.get("layout", "intro")
        model_index = MODEL_INDEX.get(layout_key, MODEL_INDEX["intro"])
        new_slide = duplicate_slide(prs, model_index)
        filler = FILLERS.get(layout_key, fill_intro)
        filler(new_slide, data)

    set_footer_everywhere(prs, presentation_title, start_index)

    # お手本スライド（元の15枚）をプレゼンテーションから除去（レイアウト／マスターはそのまま残る）
    xml_slides = prs.slides._sldIdLst
    for sld_id in original_sld_ids:
        xml_slides.remove(sld_id)

    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

# ============================================================
# Streamlit UI
# ============================================================
st.title("🎨 Slide Generator AI PRO")
st.caption("テーマや資料を入れるだけで、プロ仕様のPowerPointテンプレートに沿ったスライドを自動生成します。")

with st.sidebar:
    st.header("設定")
    num_slides = st.slider("スライド枚数（表紙・まとめ含む）", min_value=5, max_value=14, value=8)
    if os.path.exists(TEMPLATE_PATH):
        st.success("✅ master_template.pptx 準備完了")
    else:
        st.error("⚠️ master_template.pptx をapp.pyと同じフォルダに置いてください")

uploaded_file = st.file_uploader("資料をアップロード（PDF / Word）任意", type=["pdf", "docx"])
theme = st.text_area("テーマや発表したい内容を入力してください：", height=140,
                      placeholder="例：PythonによるWebスクレイピングの基礎と注意点")

if st.button("スライドを生成する", type="primary"):
    context_text = ""
    if uploaded_file:
        with st.spinner("ファイルを解析中..."):
            context_text = extract_text(uploaded_file)

    source = (context_text + "\n" + theme).strip() if context_text else theme.strip()

    if source:
        with st.spinner("AIが構成を考えています…"):
            prompt = f"""
あなたはプロのプレゼンテーションデザイナーです。以下の【内容】を元に、合計{num_slides}枚のスライド構成データを
1つのJSONオブジェクトとして出力してください。余計な説明やMarkdownの```json囲みは不要、生のJSONのみを返してください。

【内容】
{source[:4000]}

【利用できるレイアウト】
- "title"：表紙。 {{ "layout":"title", "title":"プレゼン全体のタイトル(15文字程度)", "subtitle":"発表者名やサブタイトル", "image_name":"none" }}
- "agenda"：目次。 {{ "layout":"agenda", "title":"Agenda", "items":["項目1","項目2","項目3","項目4","項目5"] }} （3〜5個）
- "chapter"：章見出し。 {{ "layout":"chapter", "eyebrow":"小さなラベル（英語可、例:'OVERVIEW'）", "title":"章のタイトル", "image_name":"none" }}
- "intro"：タイトル＋本文1段落＋画像の説明スライド。 {{ "layout":"intro", "title":"見出し(15文字程度)", "body":"60〜100文字程度の具体的な説明文", "image_name":"none" }}
- "two_content"：2カラム比較。 {{ "layout":"two_content", "title":"見出し", "columns":[{{"heading":"見出しA(10文字程度)","body":"説明文30〜50文字"}}, {{"heading":"見出しB","body":"説明文30〜50文字"}}], "image_name":"none" }}
- "three_content"：3カラムのポイント紹介。 {{ "layout":"three_content", "title":"見出し", "items":[{{"heading":"見出し1(10文字程度)","body":"説明文30〜50文字","image_name":"none"}}, ... 合計3個] }}
- "quote"：引用や強調フレーズ。 {{ "layout":"quote", "quote":"引用文や強調したい一文(40文字以内)", "author":"出典や引用元(無ければ空文字)" }}
- "summary"：まとめ。 {{ "layout":"summary", "title":"まとめ", "bullets":["要点1","要点2","要点3"], "image_name":"none" }}
- "thank_you"：最後のスライド。 {{ "layout":"thank_you", "contact_lines":["発表者名","連絡先など（無ければ省略可）"] }}

【出力フォーマット】
{{
  "presentation_title": "プレゼン全体のタイトル",
  "slides": [ 上記レイアウトのオブジェクトを順番に並べた配列 ]
}}

【構成ルール】
- 1枚目は必ず layout: "title"。2枚目は layout: "agenda"。
- 本文は "intro" "two_content" "three_content" を内容に応じて使い分け、同じレイアウトを連続させない。
- 章が複数ある場合、各章の最初に layout: "chapter" を入れる。
- 内容に印象的な一言や引用できるフレーズがあれば layout: "quote" を1枚挿入する。
- 最後から2枚目は必ず layout: "summary"、最後は必ず layout: "thank_you"。
- 合計スライド数は{num_slides}枚に近づけること。
- すべての文章は具体的な事実・固有名詞を含めること。「魅力がある」のような抽象的な表現は禁止。
- image_name は常に "none" のままでよい（画像フォルダ未整備のため）。

純粋なJSONのみを出力してください。
"""
            try:
                json_res = generate_with_retry(prompt)
                payload = json.loads(json_res)

                ppt_file = create_pptx(payload)

                st.success("🎉 スライドの生成に成功しました！下のボタンからダウンロードしてください。")
                st.download_button(
                    label="📥 パワーポイントファイルをダウンロード",
                    data=ppt_file,
                    file_name="generated_slides.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            except Exception as e:
                st.error("エラーが発生しました。もう一度お試しください。")
                st.write(e)
    else:
        st.warning("テーマを入力するか、資料をアップロードしてください。")
